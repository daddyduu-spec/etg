"""
ETG - Rappresentazione geometrica
Genera presentazione .pptx con i diagrammi fondamentali
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colori
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(50, 50, 50)
MED_GRAY = RGBColor(120, 120, 120)
LIGHT_GRAY = RGBColor(200, 200, 200)
RED = RGBColor(180, 40, 40)
BLUE = RGBColor(40, 60, 160)
GREEN = RGBColor(30, 120, 60)
ORANGE = RGBColor(200, 100, 20)
PURPLE = RGBColor(100, 40, 140)
TEAL = RGBColor(0, 128, 128)
LIGHT_BLUE = RGBColor(180, 210, 240)
LIGHT_RED = RGBColor(240, 200, 200)
LIGHT_GREEN = RGBColor(200, 235, 200)
LIGHT_YELLOW = RGBColor(255, 245, 200)
LIGHT_PURPLE = RGBColor(220, 200, 240)

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Consolas"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=BLACK, line_width=Pt(1.5)):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.color.rgb = line_color
    shape.line.width = line_width
    return shape

def add_line(slide, x1, y1, x2, y2, color=BLACK, width=Pt(1.5)):
    connector = slide.shapes.add_connector(
        1,  # straight
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def add_arrow(slide, x1, y1, x2, y2, color=BLACK, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1,
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = width
    # Freccia alla fine
    connector.end_style = 2  # arrow
    return connector

# ============================================================
# SLIDE 1: Schema generale ETG - S → U → L attraverso D
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
slide1.background.fill.solid()
slide1.background.fill.fore_color.rgb = WHITE

add_text_box(slide1, 0.5, 0.3, 12, 0.6, "ETG — Schema Generale: S → C → L", font_size=28, bold=True, color=DARK_GRAY, font_name="Arial")
add_text_box(slide1, 0.5, 0.85, 12, 0.4, "La pressione pre-grammaticale S viene ridotta da U attraverso piani disciplinari D, producendo linguaggio L", font_size=14, color=MED_GRAY, font_name="Arial")

# S a SINISTRA (convenzione Carlo)
# S grande
add_shape(slide1, MSO_SHAPE.OVAL, 0.8, 2.5, 2.2, 2.2, fill_color=LIGHT_RED, line_color=RED, line_width=Pt(3))
add_text_box(slide1, 0.8, 3.15, 2.2, 0.8, "S", font_size=48, bold=True, color=RED, alignment=PP_ALIGN.CENTER, font_name="Arial")
add_text_box(slide1, 0.5, 4.9, 2.8, 0.6, "pressione\npre-grammaticale", font_size=11, color=RED, alignment=PP_ALIGN.CENTER, font_name="Arial")

# Freccia S -> C
add_arrow(slide1, 3.1, 3.6, 4.3, 3.6, color=DARK_GRAY, width=Pt(3))

# C (la costante) - linea verticale
add_line(slide1, 4.8, 1.5, 4.8, 5.8, color=BLUE, width=Pt(4))
add_text_box(slide1, 4.4, 1.1, 0.8, 0.4, "C", font_size=32, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER, font_name="Arial")
add_text_box(slide1, 3.9, 5.9, 1.8, 0.5, "costante strutturale\n(da alfa)", font_size=10, color=BLUE, alignment=PP_ALIGN.CENTER, font_name="Arial")

# Cer (inter) - lato destro di C
add_text_box(slide1, 5.1, 1.5, 1.0, 0.3, "Cer", font_size=16, bold=True, color=GREEN, font_name="Arial")
add_text_box(slide1, 5.1, 1.8, 1.5, 0.3, "(inter)", font_size=11, color=GREEN, font_name="Arial")

# C- (intra) - lato sinistro di C
add_text_box(slide1, 3.3, 1.5, 1.0, 0.3, "C\u207b", font_size=16, bold=True, color=PURPLE, font_name="Arial")
add_text_box(slide1, 3.0, 1.8, 1.5, 0.3, "(intra)", font_size=11, color=PURPLE, font_name="Arial")

# U su C
add_shape(slide1, MSO_SHAPE.OVAL, 4.45, 2.8, 0.7, 0.7, fill_color=LIGHT_PURPLE, line_color=PURPLE, line_width=Pt(2))
add_text_box(slide1, 4.45, 2.85, 0.7, 0.6, "U", font_size=20, bold=True, color=PURPLE, alignment=PP_ALIGN.CENTER, font_name="Arial")

add_shape(slide1, MSO_SHAPE.OVAL, 4.45, 3.8, 0.7, 0.7, fill_color=LIGHT_PURPLE, line_color=PURPLE, line_width=Pt(2))
add_text_box(slide1, 4.45, 3.85, 0.7, 0.6, "U", font_size=20, bold=True, color=PURPLE, alignment=PP_ALIGN.CENTER, font_name="Arial")

# CU (dove U si sovrappongono)
add_text_box(slide1, 5.3, 3.2, 2.5, 0.4, "CU", font_size=18, bold=True, color=ORANGE, font_name="Arial")
add_text_box(slide1, 5.3, 3.55, 2.8, 0.4, "punto su C dove pi\u00f9 U coesistono", font_size=10, color=ORANGE, font_name="Arial")

# Piani D a destra (inter)
# D1
rect1 = add_shape(slide1, MSO_SHAPE.RECTANGLE, 7.0, 2.0, 2.5, 1.2, fill_color=LIGHT_BLUE, line_color=BLUE, line_width=Pt(2))
add_text_box(slide1, 7.1, 2.05, 2.3, 0.4, "D\u2081 (es. fisica)", font_size=13, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER, font_name="Arial")
add_text_box(slide1, 7.1, 2.4, 2.3, 0.7, "V\u2081, A\u2081, DH\u2081", font_size=11, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, font_name="Arial")

# D2
rect2 = add_shape(slide1, MSO_SHAPE.RECTANGLE, 7.0, 3.5, 2.5, 1.2, fill_color=LIGHT_GREEN, line_color=GREEN, line_width=Pt(2))
add_text_box(slide1, 7.1, 3.55, 2.3, 0.4, "D\u2082 (es. cucina)", font_size=13, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER, font_name="Arial")
add_text_box(slide1, 7.1, 3.9, 2.3, 0.7, "V\u2082, A\u2082, DH\u2082", font_size=11, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, font_name="Arial")

# D3
rect3 = add_shape(slide1, MSO_SHAPE.RECTANGLE, 7.0, 5.0, 2.5, 1.2, fill_color=LIGHT_YELLOW, line_color=ORANGE, line_width=Pt(2))
add_text_box(slide1, 7.1, 5.05, 2.3, 0.4, "D\u2083 (es. diritto)", font_size=13, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER, font_name="Arial")
add_text_box(slide1, 7.1, 5.4, 2.3, 0.7, "V\u2083, A\u2083, DH\u2083", font_size=11, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, font_name="Arial")

# Frecce C -> D
add_arrow(slide1, 5.2, 3.0, 6.9, 2.6, color=BLUE, width=Pt(1.5))
add_arrow(slide1, 5.2, 3.6, 6.9, 4.1, color=GREEN, width=Pt(1.5))
add_arrow(slide1, 5.2, 4.2, 6.9, 5.6, color=ORANGE, width=Pt(1.5))

# L a destra (emerge da D)
add_arrow(slide1, 9.6, 2.6, 10.5, 3.4, color=DARK_GRAY, width=Pt(1.5))
add_arrow(slide1, 9.6, 4.1, 10.5, 3.6, color=DARK_GRAY, width=Pt(1.5))
add_arrow(slide1, 9.6, 5.6, 10.5, 3.8, color=DARK_GRAY, width=Pt(1.5))

add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE, 10.5, 3.0, 2.2, 1.2, fill_color=LIGHT_GREEN, line_color=TEAL, line_width=Pt(3))
add_text_box(slide1, 10.5, 3.1, 2.2, 0.5, "L", font_size=36, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER, font_name="Arial")
add_text_box(slide1, 10.5, 3.6, 2.2, 0.5, "dominio di\nstabilizzazione", font_size=10, color=TEAL, alignment=PP_ALIGN.CENTER, font_name="Arial")

# Note in basso
add_text_box(slide1, 0.5, 6.5, 12, 0.8,
    "S a SINISTRA (convenzione) | C = costante strutturale (da alfa) | Cer = inter | C\u207b = intra | CU = punto su C, nessuna formula | L emerge da \u2113 compatibili",
    font_size=11, color=MED_GRAY, font_name="Arial")


# ============================================================
# SLIDE 2: Rettangolo D
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
slide2.background.fill.solid()
slide2.background.fill.fore_color.rgb = WHITE

add_text_box(slide2, 0.5, 0.3, 12, 0.6, "Il Rettangolo D", font_size=28, bold=True, color=DARK_GRAY, font_name="Arial")
add_text_box(slide2, 0.5, 0.85, 12, 0.4, "Ogni piano disciplinare D ha geometria rettangolare con 4 vertici e altezza DH", font_size=14, color=MED_GRAY, font_name="Arial")

# Rettangolo grande al centro
rx, ry = 2.5, 1.8  # top-left
rw, rh = 5.5, 4.0  # width, height

# Rettangolo
rect = add_shape(slide2, MSO_SHAPE.RECTANGLE, rx, ry, rw, rh, fill_color=RGBColor(245, 248, 255), line_color=BLUE, line_width=Pt(3))

# Vertici
# Top-left = V_l
add_shape(slide2, MSO_SHAPE.OVAL, rx - 0.15, ry - 0.15, 0.3, 0.3, fill_color=RED, line_color=RED)
add_text_box(slide2, rx - 0.8, ry - 0.5, 1.0, 0.4, "V_l", font_size=20, bold=True, color=RED, alignment=PP_ALIGN.CENTER, font_name="Consolas")

# Top-right = A_l
add_shape(slide2, MSO_SHAPE.OVAL, rx + rw - 0.15, ry - 0.15, 0.3, 0.3, fill_color=GREEN, line_color=GREEN)
add_text_box(slide2, rx + rw - 0.2, ry - 0.5, 1.0, 0.4, "A_l", font_size=20, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER, font_name="Consolas")

# Bottom-left = V_c
add_shape(slide2, MSO_SHAPE.OVAL, rx - 0.15, ry + rh - 0.15, 0.3, 0.3, fill_color=RED, line_color=RED)
add_text_box(slide2, rx - 0.8, ry + rh + 0.1, 1.0, 0.4, "V_c", font_size=20, bold=True, color=RED, alignment=PP_ALIGN.CENTER, font_name="Consolas")

# Bottom-right = A_c
add_shape(slide2, MSO_SHAPE.OVAL, rx + rw - 0.15, ry + rh - 0.15, 0.3, 0.3, fill_color=GREEN, line_color=GREEN)
add_text_box(slide2, rx + rw - 0.2, ry + rh + 0.1, 1.0, 0.4, "A_c", font_size=20, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER, font_name="Consolas")

# Etichette lati
# Top = L extremes
add_text_box(slide2, rx + 1.5, ry - 0.55, 3.0, 0.4, "\u2190 estremi di L \u2192", font_size=14, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER, font_name="Arial")

# Bottom = L proiettata su C
add_text_box(slide2, rx + 1.0, ry + rh + 0.15, 3.5, 0.4, "\u2190 L proiettata su C \u2192", font_size=14, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER, font_name="Arial")

# Left = S / portata
add_text_box(slide2, rx - 2.2, ry + rh/2 - 0.3, 2.0, 0.6, "S / portata\n(sinistra)", font_size=14, bold=True, color=RED, alignment=PP_ALIGN.CENTER, font_name="Arial")

# Right = attuabilita
add_text_box(slide2, rx + rw + 0.5, ry + rh/2 - 0.3, 2.0, 0.6, "attuabilit\u00e0\n(destra)", font_size=14, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER, font_name="Arial")

# DH al centro come altezza
add_arrow(slide2, rx + rw + 0.2, ry + 0.3, rx + rw + 0.2, ry + rh - 0.3, color=PURPLE, width=Pt(2.5))
add_arrow(slide2, rx + rw + 0.2, ry + rh - 0.3, rx + rw + 0.2, ry + 0.3, color=PURPLE, width=Pt(2.5))
add_text_box(slide2, rx + rw + 0.35, ry + rh/2 - 0.3, 1.5, 0.6, "DH\n(fasce di\nresistenza)", font_size=13, bold=True, color=PURPLE, alignment=PP_ALIGN.LEFT, font_name="Arial")

# D al centro
add_text_box(slide2, rx + rw/2 - 0.5, ry + rh/2 - 0.3, 1.0, 0.6, "D", font_size=48, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER, font_name="Arial")

# L tratteggiata (mediana orizzontale)
# Linea tratteggiata a meta altezza
ly = ry + rh * 0.45
for i in range(22):
    x_start = rx + 0.1 + i * 0.25
    if x_start + 0.15 < rx + rw:
        add_line(slide2, x_start, ly, x_start + 0.15, ly, color=TEAL, width=Pt(2))

add_text_box(slide2, rx + rw/2 - 1.0, ly - 0.45, 2.0, 0.35, "L (mediana)", font_size=13, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER, font_name="Arial")

# Sotto L = deterministico
add_text_box(slide2, rx + rw/2 - 1.2, ly + 0.6, 2.4, 0.35, "sotto L = deterministico", font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER, font_name="Arial")

# Sopra L = indeterministico
add_text_box(slide2, rx + rw/2 - 1.2, ly - 1.0, 2.4, 0.35, "sopra L = indeterministico", font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER, font_name="Arial")

# Legenda in basso
add_text_box(slide2, 0.5, 6.3, 12, 0.4,
    "V = operatore di valore (significativit\u00e0) | A = operatore di azione (trasformazioni) | DH = fasce di resistenza (non metrico)",
    font_size=11, color=MED_GRAY, font_name="Arial")
add_text_box(slide2, 0.5, 6.65, 12, 0.4,
    "L funge da mediana: sotto = deterministico, sopra = indeterministico | V, A, DH, \u0394 si applicano a \u2113, MAI a L",
    font_size=11, color=MED_GRAY, font_name="Arial")


# ============================================================
# SLIDE 3: ℓ - unita minima
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
slide3.background.fill.solid()
slide3.background.fill.fore_color.rgb = WHITE

add_text_box(slide3, 0.5, 0.3, 12, 0.6, "\u2113 \u2014 Unit\u00e0 Minima di Stabilizzazione Valutabile", font_size=28, bold=True, color=DARK_GRAY, font_name="Arial")
add_text_box(slide3, 0.5, 0.85, 12, 0.4, "\u2113 := \u27e8D, V(D), A(D), DH_\u2113, \u0394_\u2113\u27e9    |    Esclusivamente inter    |    Il carico \u00e8 posizionale, non nella quintupla", font_size=14, color=MED_GRAY, font_name="Arial")

# Quintupla come 5 scatole
labels = ["D", "V(D)", "A(D)", "DH_\u2113", "\u0394_\u2113"]
descriptions = [
    "Piano\ndisciplinare",
    "Operatore\ndi valore\n(significativit\u00e0)",
    "Operatore\ndi azione\n(trasformazioni)",
    "Fascia di\nresistenza",
    "Scarto\n(esito)"
]
colors_fill = [LIGHT_BLUE, LIGHT_RED, LIGHT_GREEN, LIGHT_PURPLE, LIGHT_YELLOW]
colors_line = [BLUE, RED, GREEN, PURPLE, ORANGE]

bx_start = 1.0
bw = 2.0
bh = 1.5
by = 2.0
gap = 0.3

for i, (label, desc, cf, cl) in enumerate(zip(labels, descriptions, colors_fill, colors_line)):
    x = bx_start + i * (bw + gap)
    add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, x, by, bw, bh, fill_color=cf, line_color=cl, line_width=Pt(2.5))
    add_text_box(slide3, x, by + 0.15, bw, 0.5, label, font_size=22, bold=True, color=cl, alignment=PP_ALIGN.CENTER, font_name="Consolas")
    add_text_box(slide3, x + 0.1, by + 0.7, bw - 0.2, 0.8, desc, font_size=11, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, font_name="Arial")

# Parentesi angolari
add_text_box(slide3, 0.4, 2.3, 0.6, 0.8, "\u27e8", font_size=60, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, font_name="Arial")
add_text_box(slide3, 12.2, 2.3, 0.6, 0.8, "\u27e9", font_size=60, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, font_name="Arial")

# Virgole tra le scatole
for i in range(4):
    x = bx_start + (i + 1) * (bw + gap) - gap/2
    add_text_box(slide3, x - 0.15, 2.6, 0.3, 0.4, ",", font_size=28, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, font_name="Arial")

# Sezione: cosa NON e' nella quintupla
add_text_box(slide3, 1.0, 4.0, 5.5, 0.4, "Cosa NON entra nella quintupla:", font_size=16, bold=True, color=RED, font_name="Arial")

items_no = [
    "C \u2014 \u00e8 condizione a priori, costante strutturale. Non \u00e8 propriet\u00e0 di \u2113",
    "Carico \u2014 \u00e8 propriet\u00e0 posizionale (densit\u00e0 relazionale in L). Visibile solo nel DB compilato",
    "Tempo \u2014 non esiste in ETG. Il carico sostituisce la durata"
]
for i, item in enumerate(items_no):
    add_text_box(slide3, 1.3, 4.5 + i * 0.4, 10, 0.4, "\u2716  " + item, font_size=12, color=RED, font_name="Arial")

# Sezione: vincoli
add_text_box(slide3, 1.0, 5.9, 5.5, 0.4, "Vincoli duri:", font_size=16, bold=True, color=PURPLE, font_name="Arial")

items_yes = [
    "\u2113 \u00e8 esclusivamente inter \u2014 VIETATO in intra (in intra: solo \u03bc, \u03bc')",
    "V, A, DH, \u0394 si applicano a \u2113, MAI a L \u2014 L cambia per emergenza",
    "\u0394 agisce su \u2113, induce variazioni in L \u2014 \u0394 non agisce su L come oggetto"
]
for i, item in enumerate(items_yes):
    add_text_box(slide3, 1.3, 6.35 + i * 0.35, 10, 0.35, "\u25cf  " + item, font_size=11, color=PURPLE, font_name="Arial")


# ============================================================
# SLIDE 4: I tre regimi (inter / intra / meta)
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
slide4.background.fill.solid()
slide4.background.fill.fore_color.rgb = WHITE

add_text_box(slide4, 0.5, 0.3, 12, 0.6, "I Tre Regimi: Inter, Intra, Meta", font_size=28, bold=True, color=DARK_GRAY, font_name="Arial")
add_text_box(slide4, 0.5, 0.85, 12, 0.4, "Il regime \u00e8 determinato dalla collocabilit\u00e0, non dall'intimit\u00e0 del contenuto", font_size=14, color=MED_GRAY, font_name="Arial")

# C come linea verticale centrale
cx = 6.5
add_line(slide4, cx, 1.5, cx, 6.5, color=BLUE, width=Pt(4))
add_text_box(slide4, cx - 0.3, 1.1, 0.6, 0.4, "C", font_size=28, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER, font_name="Arial")

# INTER (destra di C) = Cer
add_shape(slide4, MSO_SHAPE.ROUNDED_RECTANGLE, 7.2, 1.8, 5.3, 4.5, fill_color=LIGHT_GREEN, line_color=GREEN, line_width=Pt(2))
add_text_box(slide4, 7.4, 1.9, 5.0, 0.5, "INTER (Cer)", font_size=22, bold=True, color=GREEN, font_name="Arial")

inter_items = [
    "\u2113 := \u27e8D, V(D), A(D), DH_\u2113, \u0394_\u2113\u27e9",
    "Linguaggio espresso, collocabile",
    "Piani D con V, A, DH, \u0394",
    "DB compilabile",
    "DH applicato a \u2113",
    "\u29bf (fissazione) registrata qui",
    "",
    "Es: diario scritto, teorema pubblicato,",
    "     legge promulgata, ricetta condivisa"
]
for i, item in enumerate(inter_items):
    add_text_box(slide4, 7.6, 2.5 + i * 0.4, 4.8, 0.4, item, font_size=12, color=DARK_GRAY, font_name="Arial")

# INTRA (sinistra di C) = C-
add_shape(slide4, MSO_SHAPE.ROUNDED_RECTANGLE, 0.8, 1.8, 5.3, 4.5, fill_color=LIGHT_PURPLE, line_color=PURPLE, line_width=Pt(2))
add_text_box(slide4, 1.0, 1.9, 5.0, 0.5, "INTRA (C\u207b)", font_size=22, bold=True, color=PURPLE, font_name="Arial")

intra_items = [
    "\u03bc, \u03bc' (rielaborazioni opache)",
    "NON collocabile, NON linguistico",
    "Nessun D, nessun V, nessun A",
    "NON compilabile",
    "DH applicato a U (non a \u2113)",
    "\u2113 VIETATO qui",
    "",
    "Es: intuizione pre-verbale, disagio",
    "     senza nome, elaborazione inconscia"
]
for i, item in enumerate(intra_items):
    add_text_box(slide4, 1.0, 2.5 + i * 0.4, 4.8, 0.4, item, font_size=12, color=DARK_GRAY, font_name="Arial")

# META (sotto, attraversa tutto)
add_shape(slide4, MSO_SHAPE.ROUNDED_RECTANGLE, 2.5, 6.0, 8.5, 1.2, fill_color=LIGHT_YELLOW, line_color=ORANGE, line_width=Pt(2))
add_text_box(slide4, 2.7, 6.05, 3.0, 0.4, "META (D_meta)", font_size=18, bold=True, color=ORANGE, font_name="Arial")
add_text_box(slide4, 2.7, 6.45, 8.0, 0.7,
    "Immanente (prima e dopo) | Monodimensionale | No DH, no L, no \u0394 | Struttura logica\n\"Meta non esiste, ma c'\u00e8\" | L'impressione di DB in meta = inganno di DB distribuito di U\u2080",
    font_size=11, color=ORANGE, font_name="Arial")

# Frecce cicliche inter <-> intra
add_arrow(slide4, 6.3, 3.0, 6.7, 3.0, color=TEAL, width=Pt(2))
add_arrow(slide4, 6.7, 3.5, 6.3, 3.5, color=TEAL, width=Pt(2))
add_text_box(slide4, 5.8, 3.7, 1.5, 0.4, "\u03bc\u2192\u03bc'\u2192\u2113'", font_size=10, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER, font_name="Arial")


# ============================================================
# SLIDE 5: Pathway sensoriale
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
slide5.background.fill.solid()
slide5.background.fill.fore_color.rgb = WHITE

add_text_box(slide5, 0.5, 0.3, 12, 0.6, "Pathway Sensoriale: inter \u2192 sensi \u2192 U\u2080 \u2192 Um \u2192 inter", font_size=28, bold=True, color=DARK_GRAY, font_name="Arial")
add_text_box(slide5, 0.5, 0.85, 12, 0.4, "Um NON riceve direttamente da inter. Tutto passa per U\u2080 (riduzione).", font_size=14, color=MED_GRAY, font_name="Arial")

# Scatole del pathway
pathway_data = [
    ("Inter\n(Cer)", 1.0, GREEN, LIGHT_GREEN),
    ("Sensi", 3.5, TEAL, RGBColor(200, 235, 235)),
    ("U\u2080\n(riduzione)", 6.0, PURPLE, LIGHT_PURPLE),
    ("Um\n(lessico)", 8.5, RED, LIGHT_RED),
    ("Inter\n(Cer)", 11.0, GREEN, LIGHT_GREEN),
]

pw_y = 2.5
pw_w = 2.0
pw_h = 1.5

for label, x, lc, fc in pathway_data:
    add_shape(slide5, MSO_SHAPE.ROUNDED_RECTANGLE, x, pw_y, pw_w, pw_h, fill_color=fc, line_color=lc, line_width=Pt(2.5))
    add_text_box(slide5, x, pw_y + 0.3, pw_w, 0.9, label, font_size=18, bold=True, color=lc, alignment=PP_ALIGN.CENTER, font_name="Arial")

# Frecce tra scatole
for i in range(4):
    x1 = pathway_data[i][1] + pw_w + 0.05
    x2 = pathway_data[i+1][1] - 0.05
    add_arrow(slide5, x1, pw_y + pw_h/2, x2, pw_y + pw_h/2, color=DARK_GRAY, width=Pt(3))

# X grande: Um NON riceve da inter
add_line(slide5, 1.5, 4.5, 11.5, 4.5, color=LIGHT_GRAY, width=Pt(1))

add_text_box(slide5, 1.0, 4.8, 12, 0.5, "VIETATO: Inter \u2192 Um direttamente", font_size=18, bold=True, color=RED, alignment=PP_ALIGN.CENTER, font_name="Arial")
add_text_box(slide5, 1.0, 5.4, 12, 0.8,
    "Um (lessico, vocabolario disciplinare) non riceve input direttamente dall'inter.\nTutto il materiale sensoriale passa prima per U\u2080 (riduzione percettiva).\nU\u2080 filtra, riduce, e solo dopo il risultato arriva a Um che lo traduce in linguaggio disciplinare.",
    font_size=13, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, font_name="Arial")


# ============================================================
# SLIDE 6: Esiti di Delta
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
slide6.background.fill.solid()
slide6.background.fill.fore_color.rgb = WHITE

add_text_box(slide6, 0.5, 0.3, 12, 0.6, "Esiti di \u0394 (Scarto)", font_size=28, bold=True, color=DARK_GRAY, font_name="Arial")
add_text_box(slide6, 0.5, 0.85, 12, 0.4, "\u0394 agisce su \u2113, non su L. L cambia per emergenza come conseguenza.", font_size=14, color=MED_GRAY, font_name="Arial")

# Schema ad albero: Delta -> 4 esiti
delta_x = 5.7
delta_y = 2.0
add_shape(slide6, MSO_SHAPE.OVAL, delta_x, delta_y, 1.5, 1.0, fill_color=LIGHT_YELLOW, line_color=ORANGE, line_width=Pt(3))
add_text_box(slide6, delta_x, delta_y + 0.15, 1.5, 0.7, "\u0394", font_size=36, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER, font_name="Arial")

esiti = [
    ("\u29bf Fissazione", "L'\u2113 regge. Entra in L.\nRegistrata su C (inter).", GREEN, LIGHT_GREEN, 0.8),
    ("NI_t Scarto\ntemporaneo", "L'\u2113 non regge ora.\nPu\u00f2 ritentare.", ORANGE, LIGHT_YELLOW, 3.8),
    ("NI_tP Scarto\npermanente", "L'\u2113 non regger\u00e0 mai\nin questo D.", RED, LIGHT_RED, 6.8),
    ("\u03b5 Attraversamento", "L'\u2113 transita in un\naltro D. Potenziale K.", PURPLE, LIGHT_PURPLE, 9.8),
]

ey = 4.0
ew = 2.8
eh = 1.8

for label, desc, lc, fc, ex in esiti:
    add_shape(slide6, MSO_SHAPE.ROUNDED_RECTANGLE, ex, ey, ew, eh, fill_color=fc, line_color=lc, line_width=Pt(2))
    add_text_box(slide6, ex + 0.1, ey + 0.1, ew - 0.2, 0.7, label, font_size=14, bold=True, color=lc, alignment=PP_ALIGN.CENTER, font_name="Arial")
    add_text_box(slide6, ex + 0.1, ey + 0.9, ew - 0.2, 0.8, desc, font_size=11, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, font_name="Arial")
    # Freccia da delta all'esito
    add_arrow(slide6, delta_x + 0.75, delta_y + 1.0, ex + ew/2, ey - 0.05, color=lc, width=Pt(1.5))

# Nota
add_text_box(slide6, 0.5, 6.3, 12, 0.4,
    "K (conoscenza inter-disciplinare) accade quando \u03b5 porta un'\u2113 da D\u2081 a D\u2082 e l\u00ec regge (\u29bf)",
    font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER, font_name="Arial")
add_text_box(slide6, 0.5, 6.65, 12, 0.4,
    "Il vero K accade in intra: \u03bc \u2192 \u03bc' \u2192 \u2113' (rientro in inter)",
    font_size=12, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER, font_name="Arial")


# ============================================================
# SLIDE 7: DB Topologico - chi sta dove
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
slide7.background.fill.solid()
slide7.background.fill.fore_color.rgb = WHITE

add_text_box(slide7, 0.5, 0.3, 12, 0.6, "Il DB Topologico \u2014 Chi Pu\u00f2 Stare Dove", font_size=28, bold=True, color=DARK_GRAY, font_name="Arial")
add_text_box(slide7, 0.5, 0.85, 12, 0.4, "ETG mostra, non confronta. Il DB rende visibile la struttura.", font_size=14, color=MED_GRAY, font_name="Arial")

# Tabella
headers = ["Simbolo", "D_inter", "D_intra", "D_meta"]
header_colors = [DARK_GRAY, GREEN, PURPLE, ORANGE]

rows_data = [
    ["\u2113", "\u2714", "\u2716", "\u2716"],
    ["\u03bc / \u03bc'", "\u2716", "\u2714", "\u2716"],
    ["V, A", "\u2714", "\u2716", "\u2716"],
    ["DH", "\u2714 (su \u2113)", "\u2714 (su U)", "\u2716"],
    ["\u0394", "\u2714 (su \u2113)", "\u2716", "\u2716"],
    ["\u29bf (fissazione)", "\u2714 (su C)", "\u2716 (effetto s\u00ec)", "\u2716"],
    ["DB compilabile", "\u2714", "\u2716", "\u2716"],
    ["L", "\u2714 (emerge)", "\u2716", "\u2716"],
]

tx, ty = 1.5, 1.8
tw_col = [2.5, 2.5, 2.5, 2.5]
th_header = 0.5
th_row = 0.45

# Header
for i, (h, hc) in enumerate(zip(headers, header_colors)):
    hx = tx + sum(tw_col[:i])
    add_shape(slide7, MSO_SHAPE.RECTANGLE, hx, ty, tw_col[i], th_header, fill_color=DARK_GRAY, line_color=DARK_GRAY)
    add_text_box(slide7, hx, ty + 0.05, tw_col[i], th_header, h, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, font_name="Arial")

# Rows
for r, row in enumerate(rows_data):
    row_y = ty + th_header + r * th_row
    bg = RGBColor(248, 248, 248) if r % 2 == 0 else WHITE
    for c, cell in enumerate(row):
        cx_pos = tx + sum(tw_col[:c])
        add_shape(slide7, MSO_SHAPE.RECTANGLE, cx_pos, row_y, tw_col[c], th_row, fill_color=bg, line_color=LIGHT_GRAY, line_width=Pt(0.5))
        cell_color = GREEN if "\u2714" in cell else (RED if "\u2716" in cell else DARK_GRAY)
        if c == 0:
            cell_color = DARK_GRAY
        add_text_box(slide7, cx_pos, row_y + 0.02, tw_col[c], th_row - 0.04, cell,
                    font_size=12, bold=(c==0), color=cell_color, alignment=PP_ALIGN.CENTER, font_name="Arial")

# Tre tipi di DB
add_text_box(slide7, 1.0, 5.8, 12, 0.4, "Tre tipi di DB (proposti, non ancora formalizzati):", font_size=16, bold=True, color=BLUE, font_name="Arial")

db_types = [
    ("DB di D", "Tutto ci\u00f2 che transita su C nel dominio D \u2014 \u2113, fallimenti, NI_t, eventi parziali", BLUE),
    ("DB di L", "Istituzionale \u2014 solo \u2113 stabilizzate. Per uso esterno, middleware, trasparenza", GREEN),
    ("L stessa", "NON \u00e8 un DB \u2014 dominio storico, contiene anche l'irrecuperabile", RED),
]
for i, (name, desc, color) in enumerate(db_types):
    add_text_box(slide7, 1.3, 6.25 + i * 0.35, 2.0, 0.35, name, font_size=13, bold=True, color=color, font_name="Arial")
    add_text_box(slide7, 3.3, 6.25 + i * 0.35, 9.0, 0.35, desc, font_size=12, color=DARK_GRAY, font_name="Arial")


# ============================================================
# Salva
# ============================================================
output_path = r"C:\Users\Utente\Documents\ai\etg\ClaudeETG\07_BRAINSTORMING\ETG_Geometria.pptx"
prs.save(output_path)
print(f"Salvato: {output_path}")
print(f"Slide: {len(prs.slides)}")
